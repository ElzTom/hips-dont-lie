"""
Build PowerPoint presentation: Impact of Malnutrition on Hip Fracture Outcomes
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT    = Path(__file__).resolve().parents[1]
FC      = ROOT / "outputs" / "full_cohort"
MV      = ROOT / "outputs" / "malnourished_vs_not"
ADJ     = ROOT / "outputs" / "adjusted"
OUT     = ROOT / "presentation"
OUT.mkdir(exist_ok=True)

# ── Colours ───────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x3A, 0x5C)
MID_BLUE   = RGBColor(0x5B, 0x8D, 0xB8)
ORANGE     = RGBColor(0xE0, 0x7B, 0x54)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GREY  = RGBColor(0x44, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helpers ───────────────────────────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def text_box(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_image(slide, path, l, t, w, h):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(l), Inches(t),
                                 Inches(w), Inches(h))

def section_header(slide, label, color=MID_BLUE):
    """Thin coloured bar at top of slide with section label."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0),
                                 Inches(13.33), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    text_box(slide, label, 0.2, 0.02, 12, 0.4, size=13, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)

def slide_title(slide, title, subtitle=None):
    text_box(slide, title, 0.4, 0.55, 12.5, 0.7, size=28, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        text_box(slide, subtitle, 0.4, 1.3, 12.5, 0.5, size=16,
                 color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.LEFT)

def bullet_box(slide, items, l, t, w, h, size=15, color=DARK_GREY, indent=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ("  \u2022  " if indent else "") + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)

def divider(slide, t, color=MID_BLUE):
    line = slide.shapes.add_shape(1, Inches(0.4), Inches(t),
                                  Inches(12.5), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BLUE)

# Accent bar
bar = s.shapes.add_shape(1, Inches(0), Inches(5.8), Inches(13.33), Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()

text_box(s, "Impact of Malnutrition on Hip Fracture Outcomes",
         0.6, 1.6, 12, 1.2, size=36, bold=True, color=WHITE)
text_box(s, "Exploring associations between nutritional status at admission and\n"
            "length of stay, 30-day mortality, and discharge destination",
         0.6, 3.1, 11, 0.9, size=18, color=RGBColor(0xCC, 0xDD, 0xEE))
text_box(s, "ANZHFR SAHMRI Datathon 2026  |  n = 97,449 patients  |  2016\u20132024",
         0.6, 4.3, 11, 0.5, size=14, color=RGBColor(0xAA, 0xBB, 0xCC))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Background & Research Question
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GREY)
section_header(s, "BACKGROUND")
text_box(s, "Background & Research Question", 0.4, 0.55, 12, 0.6,
         size=24, bold=True, color=DARK_BLUE)
divider(s, 1.25)

bullet_box(s, [
    "Hip fracture is a major cause of morbidity and mortality in older adults",
    "Malnutrition is common in older hospitalised patients and may worsen outcomes",
    "The Australian & New Zealand Hip Fracture Registry (ANZHFR) collects rich clinical data across both countries",
], 0.5, 1.4, 7.5, 2.2, size=15, color=DARK_GREY)

# Research question box
rq = s.shapes.add_shape(1, Inches(0.4), Inches(3.8), Inches(12.4), Inches(1.4))
rq.fill.solid(); rq.fill.fore_color.rgb = DARK_BLUE; rq.line.fill.background()
text_box(s, "Research Question", 0.6, 3.85, 12, 0.35, size=13, bold=True, color=ORANGE)
text_box(s, "Does malnutrition at admission independently predict worse outcomes "
            "(30-day mortality, length of stay, discharge destination) "
            "in hip fracture patients?",
         0.6, 4.2, 12, 0.85, size=15, color=WHITE)

bullet_box(s, [
    "Outcomes: 30-day mortality  |  Acute ward length of stay  |  Discharge destination",
    "Exposure: Clinical malnutrition assessment (malnourished vs not malnourished)",
], 0.5, 5.4, 12, 1.0, size=13, color=DARK_GREY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Full cohort overview
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GREY)
section_header(s, "FULL COHORT — OVERVIEW")
text_box(s, "Who are these patients?", 0.4, 0.55, 12, 0.6,
         size=24, bold=True, color=DARK_BLUE)
divider(s, 1.25)

# Stat boxes
def stat_box(slide, label, value, l, t):
    box = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(2.8), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = DARK_BLUE; box.line.fill.background()
    text_box(slide, value, l+0.1, t+0.05, 2.6, 0.85, size=28, bold=True, color=ORANGE,
             align=PP_ALIGN.CENTER)
    text_box(slide, label, l+0.1, t+0.9, 2.6, 0.55, size=12, color=WHITE,
             align=PP_ALIGN.CENTER)

stat_box(s, "Total patients", "97,449",   0.4,  1.4)
stat_box(s, "Years of data",  "2016–2024", 3.5,  1.4)
stat_box(s, "Female",         "66.3%",    6.6,  1.4)
stat_box(s, "Median age",     "84 years", 9.7,  1.4)

add_image(s, FC / "03_admissions_per_year.png", 0.4, 3.1, 6.5, 4.1)
add_image(s, FC / "01_sex_distribution.png",    7.1, 3.1, 6.0, 4.1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Full cohort clinical profile
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GREY)
section_header(s, "FULL COHORT — CLINICAL PROFILE")
text_box(s, "Clinical characteristics", 0.4, 0.55, 12, 0.6,
         size=24, bold=True, color=DARK_BLUE)
divider(s, 1.25)

add_image(s, FC / "04_fracture_type.png",    0.3,  1.35, 6.3, 3.0)
add_image(s, FC / "05_asa_grade.png",        6.8,  1.35, 6.3, 3.0)
add_image(s, FC / "06_walking_ability.png",  0.3,  4.45, 6.3, 3.0)
add_image(s, FC / "10_los_distribution.png", 6.8,  4.45, 6.3, 3.0)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Malnutrition cohort
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GREY)
section_header(s, "MALNUTRITION — ASSESSED COHORT")
text_box(s, "Malnutrition assessment — who was assessed?", 0.4, 0.55, 12, 0.6,
         size=24, bold=True, color=DARK_BLUE)
divider(s, 1.25)

stat_box(s, "Patients assessed",    "50,354",  0.4, 1.4)
stat_box(s, "Malnourished",         "31.7%",   3.5, 1.4)
stat_box(s, "Not malnourished",     "68.3%",   6.6, 1.4)
stat_box(s, "Assessment rate",      "54.1%",   9.7, 1.4)

bullet_box(s, [
    "Only 54% of the full cohort had a malnutrition assessment documented",
    "Analysis restricted to assessed patients (n = 50,354)",
    "\"Not assessed\" treated as missing — not comparable to a clinical finding",
], 0.5, 3.15, 12.3, 1.1, size=14, color=DARK_GREY)

add_image(s, MV / "07_malnutrition_by_sex.png", 0.4, 4.3, 12.4, 3.0)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Baseline differences (Table 1 highlights)
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GREY)
section_header(s, "MALNOURISHED vs NOT MALNOURISHED — BASELINE")
text_box(s, "Groups differ significantly at baseline", 0.4, 0.55, 12, 0.6,
         size=24, bold=True, color=DARK_BLUE)
divider(s, 1.25)

add_image(s, MV / "05_asa_grade.png",      0.3,  1.35, 6.3, 2.9)
add_image(s, MV / "06_walking_ability.png", 6.8, 1.35, 6.3, 2.9)
add_image(s, MV / "01_age_distribution.png", 0.3, 4.35, 6.3, 2.9)
add_image(s, MV / "04_discharge_destination.png", 6.8, 4.35, 6.3, 2.9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Unadjusted outcomes
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GREY)
section_header(s, "UNADJUSTED OUTCOMES")
text_box(s, "Raw outcomes look similar — but groups are not comparable",
         0.4, 0.55, 12, 0.6, size=24, bold=True, color=DARK_BLUE)
divider(s, 1.25)

add_image(s, MV / "03_mortality_30d.png",   0.4,  1.4, 5.5, 3.5)
add_image(s, MV / "02_los_distribution.png", 6.5, 1.4, 6.5, 3.5)

bullet_box(s, [
    "30-day mortality: 7.1% (malnourished) vs 7.3% (not malnourished) — almost identical",
    "Median LOS: 7 days in both groups — no raw difference",
    "But malnourished patients are older, sicker, and frailer — making direct comparison misleading",
    "Adjusted analysis needed to isolate the true effect of malnutrition",
], 0.5, 5.1, 12.3, 2.2, size=14, color=DARK_GREY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Adjusted analysis
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GREY)
section_header(s, "ADJUSTED ANALYSIS")
text_box(s, "After adjusting for confounders — no independent effect",
         0.4, 0.55, 12, 0.6, size=24, bold=True, color=DARK_BLUE)
divider(s, 1.25)

add_image(s, ADJ / "forest_plot.png", 0.4, 1.4, 12.5, 3.6)

text_box(s, "Adjusted for:", 0.5, 5.15, 12, 0.35, size=13, bold=True, color=DARK_BLUE)
bullet_box(s, [
    "Age  |  Sex  |  ASA grade  |  Pre-admission walking ability  |  Cognitive status  |  Usual residence  |  Fracture type  |  Surgical repair"
], 0.5, 5.5, 12.3, 0.5, size=13, color=DARK_GREY, indent=False)

# Results table
rows = [
    ("Outcome", "Estimate", "95% CI", "p-value", True),
    ("30-day mortality",        "OR 0.96", "0.88 – 1.03", "0.262",  False),
    ("Discharge to aged care",  "OR 1.02", "0.94 – 1.09", "0.675",  False),
    ("Acute LOS",               "-0.1%",   "-1.2% – +1.0%","0.814", False),
]
col_x = [0.5, 5.5, 8.2, 10.8]
col_w = [4.8, 2.5, 2.4, 2.0]
row_h = 0.38
t0 = 6.1
for ri, (c1, c2, c3, c4, header) in enumerate(rows):
    t = t0 + ri * row_h
    row_color = DARK_BLUE if header else (LIGHT_GREY if ri % 2 == 0 else WHITE)
    txt_color = WHITE if header else DARK_GREY
    box = s.shapes.add_shape(1, Inches(0.4), Inches(t), Inches(12.4), Inches(row_h))
    box.fill.solid(); box.fill.fore_color.rgb = row_color; box.line.fill.background()
    for ci, (cell, cx, cw) in enumerate(zip([c1,c2,c3,c4], col_x, col_w)):
        text_box(s, cell, cx, t+0.04, cw, row_h-0.05,
                 size=12, bold=header, color=txt_color)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Longer-term mortality
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GREY)
section_header(s, "LONGER-TERM MORTALITY — 30d, 90d, 120d, 365d")
text_box(s, "No effect of malnutrition at any time point",
         0.4, 0.55, 12, 0.6, size=24, bold=True, color=DARK_BLUE)
divider(s, 1.25)

add_image(s, MV / "08_mortality_all_timepoints_unadj.png", 0.3, 1.35, 7.0, 3.5)
add_image(s, MV / "09_mortality_forest_all_timepoints.png", 7.2, 1.35, 5.9, 3.5)

# Adjusted results table
lt_rows = [
    ("Time point", "OR", "95% CI", "p-value", True),
    ("30-day",  "0.96", "0.88 – 1.03", "0.262", False),
    ("90-day",  "0.98", "0.93 – 1.05", "0.622", False),
    ("120-day", "1.01", "0.95 – 1.07", "0.824", False),
    ("365-day", "0.97", "0.92 – 1.02", "0.190", False),
]
col_x = [0.5, 4.5, 6.8, 9.8]
col_w = [3.8, 2.0, 2.8, 2.5]
row_h = 0.37
t0 = 5.1
for ri, (c1, c2, c3, c4, header) in enumerate(lt_rows):
    t = t0 + ri * row_h
    row_color = DARK_BLUE if header else (LIGHT_GREY if ri % 2 == 0 else WHITE)
    txt_color = WHITE if header else DARK_GREY
    box = s.shapes.add_shape(1, Inches(0.4), Inches(t), Inches(12.4), Inches(row_h))
    box.fill.solid(); box.fill.fore_color.rgb = row_color; box.line.fill.background()
    for cell, cx, cw in zip([c1, c2, c3, c4], col_x, col_w):
        text_box(s, cell, cx, t+0.04, cw, row_h-0.05,
                 size=12, bold=header, color=txt_color)

text_box(s,
    "Adjusted for: age, sex, ASA grade, walking ability, cognition, residence, fracture type, surgery",
    0.5, 7.05, 12, 0.35, size=11, color=DARK_GREY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Subgroup analysis
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GREY)
section_header(s, "SUBGROUP ANALYSIS — AGE, RESIDENCE, COGNITION")
text_box(s, "Consistent across every subgroup — no independent effect",
         0.4, 0.55, 12, 0.6, size=24, bold=True, color=DARK_BLUE)
divider(s, 1.25)

add_image(s, MV / "10_subgroup_mortality_forest.png", 0.3, 1.35, 7.2, 4.0)
add_image(s, MV / "11_subgroup_mortality_unadj.png",  7.4, 1.35, 5.7, 4.0)

bullet_box(s, [
    "Age < 80: OR 0.87 (p=0.068)  |  Age 80+: OR 0.99 (p=0.907)",
    "Community-dwelling: OR 0.96 (p=0.456)  |  Aged care resident: OR 0.93 (p=0.310)",
    "Cognitively intact: OR 0.98 (p=0.771)  |  Cognitively impaired: OR 0.92 (p=0.195)",
    "No subgroup shows a statistically significant effect — finding is robust",
], 0.4, 5.5, 12.5, 1.8, size=13, color=DARK_GREY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Assessment bias
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GREY)
section_header(s, "WHO GETS ASSESSED FOR MALNUTRITION?")
text_box(s, "Assessed vs not assessed — similar patients, documentation gap",
         0.4, 0.55, 12, 0.6, size=24, bold=True, color=DARK_BLUE)
divider(s, 1.25)

add_image(s, MV / "12_assessment_bias_characteristics.png", 0.3, 1.35, 7.8, 3.5)
add_image(s, MV / "13_assessment_bias_outcomes.png",        8.2, 1.35, 4.9, 3.5)

# Stats boxes
def mini_stat(slide, label, val_a, val_b, l, t):
    box = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(3.9), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = DARK_BLUE; box.line.fill.background()
    text_box(slide, f"Assessed: {val_a}  |  Not assessed: {val_b}",
             l+0.1, t+0.05, 3.7, 0.35, size=11, color=ORANGE)
    text_box(slide, label, l+0.1, t+0.42, 3.7, 0.35, size=11, color=WHITE)

mini_stat(s, "Median age",          "83 yrs",  "84 yrs",  0.4, 5.05)
mini_stat(s, "30-day mortality",    "7.2%",    "7.4%",    4.5, 5.05)
mini_stat(s, "% ASA 3-5 (sicker)", "77.7%",   "69.7%",   8.6, 5.05)

bullet_box(s, [
    "46% of patients (42,767) had no malnutrition assessment documented",
    "Assessed and unassessed groups are very similar — same age, sex, mortality, LOS",
    "Sicker patients (ASA 3-5) are MORE likely to be assessed — clinically appropriate",
    "Assessment gap likely reflects documentation/resource variation across hospitals, not patient selection",
], 0.4, 6.1, 12.5, 1.3, size=13, color=DARK_GREY)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Interpretation
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BLUE)
section_header(s, "INTERPRETATION & CONCLUSIONS", color=ORANGE)
text_box(s, "What does this mean?", 0.4, 0.55, 12, 0.6,
         size=28, bold=True, color=WHITE)
divider(s, 1.25, color=ORANGE)

points = [
    ("Malnutrition is a marker, not a driver",
     "Once baseline frailty, age, and ASA grade are accounted for, malnutrition does not independently worsen outcomes"),
    ("Groups differ at baseline",
     "Malnourished patients are older (85 vs 82), more cognitively impaired (48% vs 31%), and less mobile — not comparable without adjustment"),
    ("Assessment gap",
     "Only 54% of patients were assessed — those assessed may be a selected, higher-risk group, potentially masking the true effect"),
    ("Hospitals may be compensating",
     "Active nutritional support and geriatric review for identified malnourished patients may neutralise the effect"),
]
for i, (title, body) in enumerate(points):
    t = 1.45 + i * 1.35
    box = s.shapes.add_shape(1, Inches(0.4), Inches(t), Inches(12.4), Inches(1.2))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E, 0x45, 0x6E)
    box.line.color.rgb = MID_BLUE
    text_box(s, title, 0.6, t+0.05, 12, 0.35, size=14, bold=True, color=ORANGE)
    text_box(s, body,  0.6, t+0.45, 12, 0.65, size=13, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Limitations & Next Steps
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_GREY)
section_header(s, "LIMITATIONS & NEXT STEPS")
text_box(s, "Limitations & Next Steps", 0.4, 0.55, 12, 0.6,
         size=24, bold=True, color=DARK_BLUE)
divider(s, 1.25)

text_box(s, "Limitations", 0.5, 1.35, 5.8, 0.4, size=16, bold=True, color=DARK_BLUE)
bullet_box(s, [
    "Only 54% of patients had malnutrition assessed — selection bias possible",
    "Malnutrition tool not standardised — clinical judgement varies by site",
    "No data on nutritional interventions provided during admission",
    "Registry data — residual confounding cannot be fully excluded",
    "Short-term outcomes only — longer-term impact may differ",
], 0.5, 1.8, 5.9, 3.0, size=13, color=DARK_GREY)

text_box(s, "Next Steps", 7.1, 1.35, 5.8, 0.4, size=16, bold=True, color=DARK_BLUE)
bullet_box(s, [
    "Subgroup analysis by age group and residential status",
    "Explore 90-day and 365-day mortality for longer-term signal",
    "Investigate whether assessment rate varies by hospital",
    "Propensity score matching for a stronger causal estimate",
    "Examine interaction between malnutrition and cognitive impairment",
], 7.1, 1.8, 5.9, 3.0, size=13, color=DARK_GREY)

divider(s, 5.1)
text_box(s,
    "Key message: In this registry cohort, malnutrition at admission is a marker of frailty "
    "rather than an independent predictor of hip fracture outcomes. "
    "Improving malnutrition assessment rates and standardising tools are important priorities.",
    0.5, 5.2, 12.3, 1.8, size=14, bold=False, color=DARK_BLUE)

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out_path = OUT / "malnutrition_hip_fracture_v4.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
